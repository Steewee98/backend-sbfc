import io
import traceback
from flask import Blueprint, request, jsonify, send_file
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

brandizzatore_bp = Blueprint('brandizzatore', __name__)

COLORE_SCURO = RGBColor(0x37, 0x39, 0x3f)
COLORE_TERRACOTTA = RGBColor(0xc4, 0x62, 0x2d)
COLORE_AVORIO = RGBColor(0xf5, 0xf2, 0xee)


@brandizzatore_bp.route('/api/brandizza/word', methods=['POST'])
def brandizza_word():
    if 'file' not in request.files:
        return jsonify({'error': 'Nessun file'}), 400

    file = request.files['file']
    doc = Document(file)

    # Modifica stili paragrafi
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading'):
            for run in para.runs:
                run.font.color.rgb = COLORE_TERRACOTTA
                run.font.bold = True
        else:
            for run in para.runs:
                if run.font.color and run.font.color.type is not None:
                    run.font.color.rgb = COLORE_SCURO

    # Modifica stili nelle tabelle
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.type is not None:
                            run.font.color.rgb = COLORE_SCURO

    # Header su tutte le sezioni
    for section in doc.sections:
        header = section.header
        header.is_linked_to_previous = False
        if header.paragraphs:
            header.paragraphs[0].clear()
        else:
            header.add_paragraph()
        header_para = header.paragraphs[0]
        header_run = header_para.add_run('SB Food Consulting')
        header_run.font.color.rgb = COLORE_TERRACOTTA
        header_run.font.bold = True
        header_run.font.size = Pt(10)
        header_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

        # Linea terracotta sotto header
        border_run = header_para.add_run('\n')
        border_run.font.size = Pt(2)

        # Footer
        footer = section.footer
        footer.is_linked_to_previous = False
        if footer.paragraphs:
            footer.paragraphs[0].clear()
        else:
            footer.add_paragraph()
        footer_para = footer.paragraphs[0]
        footer_run = footer_para.add_run(
            'SB Food Consulting  |  info@sbfoodconsulting.com  |  www.sbfoodconsulting.com'
        )
        footer_run.font.color.rgb = COLORE_SCURO
        footer_run.font.size = Pt(8)
        footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    output = io.BytesIO()
    doc.save(output)
    output.seek(0)

    return send_file(
        output,
        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        as_attachment=True,
        download_name='contratto-sbfood-brandizzato.docx'
    )


@brandizzatore_bp.route('/api/brandizza/pptx', methods=['POST'])
def brandizza_pptx():
    if 'file' not in request.files:
        return jsonify({'error': 'Nessun file'}), 400

    try:
        file = request.files['file']
        prs = Presentation(file)

        SCURO = PptxRGB(0x37, 0x39, 0x3f)
        TERRACOTTA = PptxRGB(0xc4, 0x62, 0x2d)
        AVORIO = PptxRGB(0xf5, 0xf2, 0xee)
        BIANCO = PptxRGB(0xFF, 0xFF, 0xFF)

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for i, slide in enumerate(prs.slides):
            # SFONDO — alterna scuro/chiaro
            fill = slide.background.fill
            fill.solid()
            if i % 2 == 0:
                fill.fore_color.rgb = SCURO
                colore_testo = BIANCO
                colore_titolo = TERRACOTTA
                colore_linea = TERRACOTTA
            else:
                fill.fore_color.rgb = AVORIO
                colore_testo = SCURO
                colore_titolo = SCURO
                colore_linea = TERRACOTTA

            # LINEA DECORATIVA sinistra
            line = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR_TYPE.STRAIGHT
                PptxInches(0.3), PptxInches(0.8),
                PptxInches(0.3), PptxInches(6.5)
            )
            line.line.color.rgb = TERRACOTTA
            line.line.width = PptxPt(2)

            # COLORA TUTTI I TESTI esistenti
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Titoli principali
                        if (shape == slide.shapes[0] and
                                run.font.size and
                                run.font.size >= PptxPt(20)):
                            run.font.color.rgb = colore_titolo
                            run.font.bold = True
                        else:
                            run.font.color.rgb = colore_testo

                        # Font moderno
                        run.font.name = 'Calibri'

            # LOGO SB FOOD CONSULTING in alto a destra
            logo_box = slide.shapes.add_textbox(
                slide_width - PptxInches(2.8),
                PptxInches(0.15),
                PptxInches(2.6),
                PptxInches(0.45)
            )
            logo_tf = logo_box.text_frame
            logo_p = logo_tf.paragraphs[0]
            logo_run = logo_p.add_run()
            logo_run.text = 'SB Food Consulting'
            logo_run.font.size = PptxPt(9)
            logo_run.font.color.rgb = TERRACOTTA
            logo_run.font.bold = True
            logo_run.font.name = 'Calibri'
            logo_p.alignment = PP_ALIGN.RIGHT

            # NUMERO SLIDE in basso a destra
            num_box = slide.shapes.add_textbox(
                slide_width - PptxInches(1.2),
                slide_height - PptxInches(0.5),
                PptxInches(1.0),
                PptxInches(0.4)
            )
            num_tf = num_box.text_frame
            num_p = num_tf.paragraphs[0]
            num_run = num_p.add_run()
            num_run.text = str(i + 1)
            num_run.font.size = PptxPt(9)
            num_run.font.color.rgb = TERRACOTTA
            num_run.font.name = 'Calibri'
            num_p.alignment = PP_ALIGN.RIGHT

            # SITO WEB in basso a sinistra
            sito_box = slide.shapes.add_textbox(
                PptxInches(0.5),
                slide_height - PptxInches(0.5),
                PptxInches(3),
                PptxInches(0.4)
            )
            sito_tf = sito_box.text_frame
            sito_p = sito_tf.paragraphs[0]
            sito_run = sito_p.add_run()
            sito_run.text = 'sbfoodconsulting.com'
            sito_run.font.size = PptxPt(8)
            sito_run.font.color.rgb = (
                TERRACOTTA if i % 2 == 0
                else SCURO
            )
            sito_run.font.name = 'Calibri'

        # SLIDE DI COPERTINA — inserisci all'inizio
        layout_index = min(6, len(prs.slide_layouts) - 1)
        slide_copertina = prs.slides.add_slide(
            prs.slide_layouts[layout_index])

        # Sposta la copertina all'inizio
        xml_slides = prs.slides._sldIdLst
        last = xml_slides[-1]
        xml_slides.remove(last)
        xml_slides.insert(0, last)

        # Sfondo copertina scuro
        fill = slide_copertina.background.fill
        fill.solid()
        fill.fore_color.rgb = SCURO

        # Linea decorativa sinistra copertina
        line_cop = slide_copertina.shapes.add_connector(
            1,
            PptxInches(0.4), PptxInches(1.5),
            PptxInches(0.4), PptxInches(5.5)
        )
        line_cop.line.color.rgb = TERRACOTTA
        line_cop.line.width = PptxPt(3)

        # Brand name
        brand_box = slide_copertina.shapes.add_textbox(
            PptxInches(0.8), PptxInches(1.5),
            PptxInches(8.5), PptxInches(0.6)
        )
        brand_tf = brand_box.text_frame
        brand_p = brand_tf.paragraphs[0]
        brand_run = brand_p.add_run()
        brand_run.text = 'SB FOOD CONSULTING'
        brand_run.font.size = PptxPt(11)
        brand_run.font.color.rgb = TERRACOTTA
        brand_run.font.bold = True
        brand_run.font.name = 'Calibri'

        # Titolo presentazione
        titolo_box = slide_copertina.shapes.add_textbox(
            PptxInches(0.8), PptxInches(2.2),
            PptxInches(8.5), PptxInches(2.5)
        )
        titolo_tf = titolo_box.text_frame
        titolo_tf.word_wrap = True
        titolo_p = titolo_tf.paragraphs[0]
        titolo_run = titolo_p.add_run()
        # Prende il titolo dalla prima slide originale
        titolo_testo = 'Presentazione'
        try:
            for shape in prs.slides[1].shapes:
                if shape.has_text_frame:
                    testo = shape.text_frame.text.strip()
                    if testo and len(testo) > 3:
                        titolo_testo = testo[:80]
                        break
        except Exception:
            pass
        titolo_run.text = titolo_testo
        titolo_run.font.size = PptxPt(32)
        titolo_run.font.color.rgb = BIANCO
        titolo_run.font.bold = True
        titolo_run.font.name = 'Calibri'

        # Anno e sito
        info_box = slide_copertina.shapes.add_textbox(
            PptxInches(0.8), PptxInches(5.2),
            PptxInches(8), PptxInches(0.6)
        )
        info_tf = info_box.text_frame
        info_p = info_tf.paragraphs[0]
        info_run = info_p.add_run()
        info_run.text = 'sbfoodconsulting.com  |  info@sbfoodconsulting.com'
        info_run.font.size = PptxPt(9)
        info_run.font.color.rgb = PptxRGB(0x9a, 0x96, 0x90)
        info_run.font.name = 'Calibri'

        # SLIDE FINALE CTA
        slide_finale = prs.slides.add_slide(
            prs.slide_layouts[layout_index])

        fill = slide_finale.background.fill
        fill.solid()
        fill.fore_color.rgb = TERRACOTTA

        # Testo CTA
        cta_box = slide_finale.shapes.add_textbox(
            PptxInches(0.8), PptxInches(2.0),
            PptxInches(8.5), PptxInches(1.5)
        )
        cta_tf = cta_box.text_frame
        cta_tf.word_wrap = True
        cta_p = cta_tf.paragraphs[0]
        cta_run = cta_p.add_run()
        cta_run.text = 'Prenota una\nconsulenza gratuita.'
        cta_run.font.size = PptxPt(36)
        cta_run.font.color.rgb = BIANCO
        cta_run.font.bold = True
        cta_run.font.name = 'Calibri'
        cta_p.alignment = PP_ALIGN.LEFT

        # Link Calendly
        link_box = slide_finale.shapes.add_textbox(
            PptxInches(0.8), PptxInches(4.0),
            PptxInches(8.5), PptxInches(0.8)
        )
        link_tf = link_box.text_frame
        link_p = link_tf.paragraphs[0]
        link_run = link_p.add_run()
        link_run.text = 'calendly.com/sbfoodconsulting-info/30min'
        link_run.font.size = PptxPt(14)
        link_run.font.color.rgb = BIANCO
        link_run.font.name = 'Calibri'

        # Brand finale
        brand_fin_box = slide_finale.shapes.add_textbox(
            PptxInches(0.8), PptxInches(5.2),
            PptxInches(8), PptxInches(0.5)
        )
        brand_fin_tf = brand_fin_box.text_frame
        brand_fin_p = brand_fin_tf.paragraphs[0]
        brand_fin_run = brand_fin_p.add_run()
        brand_fin_run.text = 'SB FOOD CONSULTING'
        brand_fin_run.font.size = PptxPt(9)
        brand_fin_run.font.color.rgb = BIANCO
        brand_fin_run.font.bold = True
        brand_fin_run.font.name = 'Calibri'

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='presentazione-sbfood-brandizzata.pptx'
        )

    except Exception as e:
        traceback.print_exc()
        print(f"ERRORE PPTX: {e}")
        return jsonify({'error': str(e)}), 500
